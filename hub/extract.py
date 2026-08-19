"""Extraction tiers: MinerU (PDF/image), native (Office/email), metadata.

Every function returns (text, tier, quality_note). A missing tool or a
crashing parser degrades to a lower tier with a visible note — ingestion
never fails because an extractor is unavailable.
"""

import email
import shutil
import subprocess
import tempfile
from datetime import datetime, timezone
from email import policy
from pathlib import Path

MINERU_EXTENSIONS = {".pdf", ".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp"}
OFFICE_EXTENSIONS = {".docx": "docx", ".pptx": "pptx", ".xlsx": "xlsx"}
EMAIL_EXTENSIONS = {".eml", ".msg"}
CAD_EXTENSIONS = {".dwg", ".dxf", ".rvt", ".rfa", ".step", ".stp", ".iges", ".igs", ".ifc", ".skp", ".3dm", ".nx", ".prt", ".sldrw", ".sldprt"}


def kind_for_extension(ext):
    ext = ext.lower()
    if ext == ".pdf":
        return "pdf"
    if ext in MINERU_EXTENSIONS:
        return "image"
    if ext in OFFICE_EXTENSIONS:
        return "office"
    if ext in EMAIL_EXTENSIONS:
        return "email"
    if ext in CAD_EXTENSIONS:
        return "cad"
    return "other"


TEXT_EXTENSIONS = {".txt", ".md"}


def extract(path_str, mineru_path="mineru", mineru_timeout=900):
    """Route a file to its tier. Returns (text, tier, quality_note)."""
    path = Path(path_str)
    ext = path.suffix.lower()

    if ext in TEXT_EXTENSIONS:
        try:
            return (
                path.read_text(encoding="utf-8", errors="replace"),
                "native",
                "",
            )
        except OSError:
            return _metadata_block(path), "metadata", "unreadable text file"
    if ext in OFFICE_EXTENSIONS:
        text = _extract_office(path, ext)
        if text is not None:
            return text, "native", ""
    if ext in EMAIL_EXTENSIONS:
        text = _extract_email(path)
        if text is not None:
            return text, "email", ""
    if ext in MINERU_EXTENSIONS:
        return _extract_pdf_or_image(path, mineru_path, mineru_timeout)
    return _metadata_block(path), "metadata", "no machine-readable content"


def mineru_available(mineru_path="mineru"):
    if Path(mineru_path).is_file():
        return True
    return shutil.which(mineru_path) is not None


def _extract_pdf_or_image(path, mineru_path, mineru_timeout):
    if mineru_available(mineru_path):
        text = _run_mineru(path, mineru_path, mineru_timeout)
        if text is not None:
            return text, "mineru", ""
    # Fallback 1: text layer via pypdf
    if path.suffix.lower() == ".pdf":
        text = _extract_pdf_text(path)
        if text and text.strip():
            return (
                text,
                "native",
                "MinerU unavailable — used fast text-layer extraction; "
                "scanned pages were skipped",
            )
        return (
            _metadata_block(path),
            "metadata",
            "scanned PDF — install MinerU for OCR extraction",
        )
    return _metadata_block(path), "metadata", "image file — install MinerU for OCR extraction"


def _run_mineru(path, mineru_path, mineru_timeout):
    with tempfile.TemporaryDirectory(prefix="rph_mineru_") as outdir:
        cmd = [mineru_path, "-p", str(path), "-o", outdir]
        try:
            proc = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=mineru_timeout,
            )
        except (subprocess.TimeoutExpired, OSError) as exc:
            return None
        if proc.returncode != 0:
            return None
        md_files = sorted(
            Path(outdir).rglob("*.md"), key=lambda p: p.stat().st_mtime
        )
        if not md_files:
            return None
        return md_files[-1].read_text(encoding="utf-8", errors="replace")


def _extract_pdf_text(path):
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        pages = []
        for page in reader.pages:
            pages.append(page.extract_text() or "")
        return "\n\n".join(pages)
    except Exception:
        return None


def _extract_office(path, ext):
    try:
        if ext == ".docx":
            return _extract_docx(path)
        if ext == ".pptx":
            return _extract_pptx(path)
        if ext == ".xlsx":
            return _extract_xlsx(path)
    except Exception:
        return None
    return None


def _extract_docx(path):
    import docx

    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for table in d.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(path):
    from pptx import Presentation

    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        chunks = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        chunks.append(text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                chunks.append(f"[notes] {notes}")
        if len(chunks) > 1:
            parts.append("\n".join(chunks))
    return "\n\n".join(parts)


def _extract_xlsx(path):
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"--- Sheet: {ws.title} ---\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


def _extract_email(path):
    if path.suffix.lower() == ".msg":
        return _extract_msg(path)
    raw = path.read_bytes()
    msg = email.message_from_bytes(raw, policy=policy.default)
    headers = [
        f"From: {msg.get('From', '')}",
        f"To: {msg.get('To', '')}",
        f"Date: {msg.get('Date', '')}",
        f"Subject: {msg.get('Subject', '')}",
    ]
    body = ""
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                body = part.get_content()
                break
    else:
        body = msg.get_content()
    return "\n".join(headers) + "\n\n" + (body or "")


def _extract_msg(path):
    import extract_msg

    m = extract_msg.Message(str(path))
    try:
        return (
            f"From: {m.sender}\nTo: {m.to}\n"
            f"Date: {m.date}\nSubject: {m.subject}\n\n{m.body}"
        )
    finally:
        m.close()


def _metadata_block(path):
    stat = path.stat()
    mtime = datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc)
    return (
        f"Filename: {path.name}\n"
        f"Size: {stat.st_size} bytes\n"
        f"Last modified: {mtime.strftime('%Y-%m-%d %H:%M UTC')}\n"
        "[no machine-readable content — binary/native CAD or unsupported format]"
    )
