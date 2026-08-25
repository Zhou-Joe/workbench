"""Shared fixture builders: tiny real files in a temp workspace."""

import tempfile
from email.message import EmailMessage
from pathlib import Path

from django.test import TestCase

from hub.models import AppSettings, Project, make_phase


def make_text_pdf(path, text):
    """A minimal single-page PDF with one line of real text."""
    stream = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode()
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R"
        b" /Resources << /Font << /F1 5 0 R >> >> >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    out = b"%PDF-1.4\n"
    offsets = []
    for i, body in enumerate(objs, 1):
        offsets.append(len(out))
        out += f"{i} 0 obj\n".encode() + body + b"\nendobj\n"
    xref_pos = len(out)
    out += f"xref\n0 {len(objs) + 1}\n0000000000 65535 f \n".encode()
    for off in offsets:
        out += f"{off:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<< /Size {len(objs) + 1} /Root 1 0 R >>\nstartxref\n{xref_pos}\n%%EOF"
    ).encode()
    Path(path).write_bytes(out)


def make_docx(path, lines):
    import docx

    d = docx.Document()
    for line in lines:
        d.add_paragraph(line)
    d.save(str(path))


def make_pptx(path, slides):
    from pptx import Presentation

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for title, body in slides:
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
        tf = box.text_frame
        tf.text = title
        p = tf.add_paragraph()
        p.text = body
    prs.save(str(path))


def make_xlsx(path, rows):
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Status"
    for row in rows:
        ws.append(row)
    wb.save(str(path))


def make_eml(path, subject, body, from_="vendor@partner.com", to="me@example.com"):
    msg = EmailMessage()
    msg["From"] = from_
    msg["To"] = to
    msg["Subject"] = subject
    msg["Date"] = "Wed, 12 Aug 2026 09:30:00 +0800"
    msg.set_content(body)
    Path(path).write_bytes(msg.as_bytes())


def make_fake_cad(path):
    Path(path).write_bytes(b"\x00" * 128 + b"FAKE-DWG-CONTENT" + b"\x00" * 128)


class WorkspaceTestCase(TestCase):
    """Test case with a temp workspace root and one seeded project."""

    def setUp(self):
        super().setUp()
        self._tmp = tempfile.TemporaryDirectory(prefix="rph_test_ws_")
        self.addCleanup(self._tmp.cleanup)
        self.root = Path(self._tmp.name)
        self.settings_obj = AppSettings.load()
        self.settings_obj.workspace_root = str(self.root)
        self.settings_obj.save()

    def seed_project(self, name="Cosmic Coaster", slug="cosmic-coaster"):
        project = Project.objects.create(name=name, slug=slug, code="CC")
        for order in range(1, 4):
            make_phase(project, f"Phase {order}", order)
        proj_dir = self.root / slug
        for order in range(1, 4):
            (proj_dir / f"{order:02d}-phase-{order}").mkdir(parents=True, exist_ok=True)
        return project

    def phase_dir(self, project, order):
        phase = project.phases.get(order=order)
        return self.root / project.slug / phase.folder_name
